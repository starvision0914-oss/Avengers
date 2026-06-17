# -*- coding: utf-8 -*-
"""지마켓·옥션 매출 증가 세부 대책 보고서 PPT 생성 (python-pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x1F, 0x33, 0x66)
ORANGE = RGBColor(0xE6, 0x77, 0x00)
RED = RGBColor(0xD3, 0x2F, 0x2F)
GREEN = RGBColor(0x1B, 0x7A, 0x46)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF2, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _tf(box):
    tf = box.text_frame; tf.word_wrap = True; return tf


def bar(slide, color=NAVY, h=Inches(0.18), top=0):
    s = slide.shapes.add_shape(1, 0, Emu(int(top)), SW, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s


def title_slide(title, subtitle, foot):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH); bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    accent = s.shapes.add_shape(1, 0, Inches(4.6), SW, Inches(0.12)); accent.fill.solid(); accent.fill.fore_color.rgb = ORANGE; accent.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.6)); tf = _tf(t)
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = WHITE
    st = s.shapes.add_textbox(Inches(0.95), Inches(4.9), Inches(11.5), Inches(1.0)); tf2 = _tf(st)
    p2 = tf2.paragraphs[0]; p2.text = subtitle; p2.font.size = Pt(18); p2.font.color.rgb = RGBColor(0xCF,0xD8,0xEC)
    f = s.shapes.add_textbox(Inches(0.95), Inches(6.7), Inches(11.5), Inches(0.5)); tf3 = _tf(f)
    p3 = tf3.paragraphs[0]; p3.text = foot; p3.font.size = Pt(12); p3.font.color.rgb = RGBColor(0x9F,0xAE,0xCE)
    return s


def header(slide, no, title):
    bar(slide, NAVY, Inches(0.85))
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.62)); tf = _tf(t)
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = f"{no}  "; r1.font.size = Pt(16); r1.font.bold = True; r1.font.color.rgb = ORANGE
    r2 = p.add_run(); r2.text = title; r2.font.size = Pt(22); r2.font.bold = True; r2.font.color.rgb = WHITE


def content_slide(no, title, bullets, note=None):
    s = prs.slides.add_slide(BLANK); header(s, no, title)
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.1), Inches(5.7)); tf = _tf(box)
    first = True
    for lvl, text, color in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.level = lvl
        run = p.add_run(); run.text = text
        run.font.size = Pt(20 - lvl*3) if lvl < 2 else Pt(14)
        run.font.color.rgb = color or (NAVY if lvl == 0 else GRAY)
        run.font.bold = (lvl == 0)
        p.space_after = Pt(6 if lvl == 0 else 3)
    if note:
        nb = s.shapes.add_shape(1, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.6))
        nb.fill.solid(); nb.fill.fore_color.rgb = LIGHT; nb.line.fill.background()
        ntf = _tf(nb); ntf.margin_top = Pt(4); ntf.margin_left = Pt(10)
        np = ntf.paragraphs[0]; np.text = note; np.font.size = Pt(12); np.font.color.rgb = GRAY; np.font.italic = True
    return s


def table_slide(no, title, headers, rows, widths=None, note=None, hl=None):
    s = prs.slides.add_slide(BLANK); header(s, no, title)
    nrows, ncols = len(rows)+1, len(headers)
    tb = s.shapes.add_table(nrows, ncols, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.4*nrows)).table
    if widths:
        for i, w in enumerate(widths): tb.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        c = tb.cell(0, j); c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        pr = c.text_frame.paragraphs[0]; pr.font.size = Pt(13); pr.font.bold = True; pr.font.color.rgb = WHITE; pr.alignment = PP_ALIGN.CENTER
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = tb.cell(i, j); c.text = str(val)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            pr = c.text_frame.paragraphs[0]; pr.font.size = Pt(12)
            pr.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            if hl and hl(i-1, row): pr.font.color.rgb = RED; pr.font.bold = True
    if note:
        nt = s.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.5)); ntf = _tf(nt)
        p = ntf.paragraphs[0]; p.text = note; p.font.size = Pt(12); p.font.italic = True; p.font.color.rgb = GRAY
    return s


# ===== 1. 표지 =====
title_slide("지마켓·옥션 매출 증가 세부 대책 보고서",
            "데이터 진단 + 지마켓 공식정책 리서치 기반 실행 해결책  ·  ESM PLUS 14계정 운영",
            "작성: 운영팀  ·  2026-06-13  ·  출처: ESM PLUS 공식 판매자가이드 외 (딥리서치 2회·검증 32클레임)")

# ===== 2. Executive Summary =====
content_slide("요약", "핵심 결론 — 매출 회복의 1순위는 '광고'가 아니라 '카탈로그 노출'", [
    (0, "광고는 이미 흑자(전체 ROAS 262.8%) — 문제는 광고 효율이 아니다", GREEN),
    (0, "주문수 붕괴가 진짜 문제: 2026-03 1,227건 → 2026-05 434건 (-65%)", RED),
    (1, "공식정책상 노출은 '판매중·노출가능' 상품만 → 품절·판매중지·노출제한·중복등록 감점이 노출 모수를 줄임", None),
    (0, "그래서 우선순위는 P0 카탈로그 정상화 → P1 가격·광고 최적화 → P2 전환율·소싱", NAVY),
    (1, "P0: 14계정 전 상품 판매·노출상태 전수점검, 판매중 복구, 중복등록 정리", None),
    (1, "P1: 적자 133개 광고 OFF(월 -50만원 낭비차단) → 고ROAS 51개 예산↑ / 무료배송·최저가", None),
    (1, "P2: 클릭0전환 상품 = 경쟁사 최저가 이탈 진단, 리뷰 확보, 도매매 소싱으로 신상품 확대", None),
], note="자원을 광고 증액이 아니라 '노출 정상화'에 먼저 투입하는 것이 핵심 메시지")

# ===== 3. 현황 진단 =====
table_slide("진단", "현황 — 매출·주문수 급감 vs 광고 흑자",
    ["월", "매출(원)", "주문수", "비고"],
    [["2026-02", "51,528,281", "909", ""],
     ["2026-03", "66,946,033", "1,227", "← 정점"],
     ["2026-04", "52,613,632", "1,084", ""],
     ["2026-05", "21,489,149", "434", "← 급감 -65%"],
     ["2026-06", "10,815,991", "235", "(13일까지)"]],
    widths=[2.5,3.5,2.5,3.8],
    note="당월 광고: 광고비 3,153,887 / 광고전환매출 8,287,340 / 전체 ROAS 262.8%(흑자). → 광고가 아니라 '팔 상품 노출'이 줄어든 것.",
    hl=lambda i,r: '급감' in r[3])

# ===== 4. 원인 규명 =====
content_slide("원인", "원인 규명 — 주문수 붕괴 = 노출 모수 감소", [
    (0, "지마켓 공식정책: 검색노출은 '관리자 확인된 판매가능' 상품에만 허용", NAVY),
    (1, "품절·판매중지·노출제한·노출대기(매매부적합) 상품은 검색에서 빠짐", None),
    (1, "발송지연 누적 → 노출제한 유발 / 중복등록·부정키워드 → 정렬점수 감점", None),
    (0, "즉 주문수 1,227→434는 '팔 상품이 검색에 안 떠서' 생긴 결과일 개연성이 큼", RED),
    (1, "우리 시스템 데이터로도 확인: 판매중지 상태가 비고에 거의 안 잡히는 등 카탈로그 신선도 이슈 존재", None),
    (0, "광고 ROAS는 흑자라, 광고는 '정상 노출' 위에서 미세조정하면 됨", GREEN),
], note="출처: ESM SellerGuide2 comparePrice.html / goodsExposeWait.html (3-0 검증)")

# ===== 5. 대책 P0 =====
content_slide("대책 P0", "[최우선] 카탈로그 정상화 — 노출 모수 복구", [
    (0, "14계정 전 상품의 판매상태·노출상태 전수 점검", NAVY),
    (1, "품절/판매중지/판매불가 → 재고확보 후 '판매중' 복구, 노출제한·노출대기 해제", None),
    (1, "발송지연 해소(노출제한 원인 제거)", None),
    (0, "중복등록 정리 (어뷰징 감점 회피)", NAVY),
    (1, "동일상품을 여러 계정·근접 카테고리에 중복 등록 시 정렬점수 감점 → 노출↓", None),
    (1, "사은품·수량만 다른 실질 동일상품도 중복으로 간주됨 → 정리", None),
    (0, "KPI: 14계정 '노출가능 상품수' 회복(03 수준), 판매중 비율 ↑", GREEN),
], note="출처: comparePrice.html, syi_policy.html, policy.gmarket.co.kr/seller-use (3-0)")

# ===== 6. 대책 P1-A 노출/가격 =====
content_slide("대책 P1-A", "노출 랭킹·가격 경쟁력 최적화", [
    (0, "가격비교 노출순위 = 최저가 > 무료배송 > 랭크(판매실적·검색정확도·고객행태·서비스품질)", NAVY),
    (1, "'최저가'는 판매가가 아니라 쿠폰적용가+배송비 합산 → 무료배송·가격인하가 노출 1·2순위 레버", None),
    (1, "적자광고 절감액(월 50만)을 일부 '무료배송 전환·가격 인하'로 돌려 노출 확보", None),
    (0, "검색용 상품명에 핵심 키워드 배치 (프로모션용은 검색 제외)", NAVY),
    (1, "검색용은 등록 10일 이내·무판매 상태에서만 수정 가능 → 신규등록 초기 키워드 확정 필수", None),
    (1, "같은 의미라도 붙여쓰기/띄어쓰기 검색량 다름 → 둘 다 분석 후 유리한 형태 선택", None),
    (1, "리뉴얼·신규는 '상품2.0'으로 등록 → 검색 랭킹 가점", None),
], note="출처: goodsName.html, productRegisterIntegrate.html, i-boss 키워드가이드")

# ===== 7. 대책 P1-B 광고운영룰 =====
content_slide("대책 P1-B", "광고 운영 룰 — 적자 OFF / 고ROAS 증액", [
    (0, "파워클릭 CPC: 최저 90원·10원 단위, 노출=입찰가 × 품질평가점수 × 구매고객특성", NAVY),
    (1, "입찰가만으로 순위 안 올라감 → 상품명·이미지·판매실적(품질점수) 개선이 곧 ROAS 개선", None),
    (0, "적자 133개(ROAS≤100, 월 50만 낭비): 광고 OFF / 입찰 축소", RED),
    (0, "고ROAS 51개(광고비 찔끔→전환 280만): 입찰·일예산 확대 + AI매출업 병행", GREEN),
    (1, "AI매출업은 입찰가까지 AI 자동최적화 → 과금 급증 위험 → 일예산·목표ROAS 상한으로 통제", None),
    (1, "품절·판매중지 상품은 광고 등록 자체 불가 → 'P0 정상화' 후에만 광고룰 성립", None),
    (0, "수치 임계치는 공식기준 없음 → 우리 자체 데이터로 도출(아래 슬라이드)", ORANGE),
], note="출처: marketinghub power-click/introduction, AD_Guide_all.pdf. 입찰가 역산 통제 주장은 반증(0-3)되어 불채택")

# ===== 8. 자체데이터 기반 룰 =====
table_slide("운영 룰", "광고 자동화 룰 — 자체 데이터 기반 수치 기준(제안)",
    ["구분", "조건(자체데이터)", "조치"],
    [["적자 OFF", "당월 광고비≥2,000 & 클릭≥10 & 전환0(또는 ROAS≤100)", "광고 OFF·입찰축소 + 상세/가격 점검"],
     ["고ROAS 증액", "ROAS≥300 & 전환>0 (특히 광고비 과소)", "입찰·일예산 단계적 ↑ (과금 모니터링)"],
     ["관찰", "ROAS 100~300", "유지·소폭 조정"],
     ["상한", "AI매출업 적용 상품", "일예산·목표ROAS 상한 필수"]],
    widths=[2.2,6.3,4.0],
    note="공식 출처상 '클릭10·전환0 OFF', '목표ROAS 500/700%' 같은 일반수치는 검증 실패 → 우리 데이터(적자133·고ROAS51)로 임계치 확정함")

# ===== 9. 대책 P2 전환율 =====
content_slide("대책 P2", "전환율(CVR) 개선 — '클릭은 나는데 전환 0' 처방", [
    (0, "온라인 평균 구매전환율 ~2% — 클릭 많은데 전환0이면 명백한 이상신호", NAVY),
    (0, "1순위 원인: 가격비교에서 최저가가 아니라 '동일상품 더 싼 경쟁셀러로 이탈'", RED),
    (1, "예) 북카트 광고비 48,664·클릭25·전환0 → 경쟁 최저가 점검 우선", None),
    (0, "처방 우선순위", NAVY),
    (1, "① 가격·배송비(쿠폰적용가+배송비) 경쟁력 확보(무료배송) ", None),
    (1, "② 리뷰 확보 — 리뷰 부재는 강력한 전환 이탈(고가일수록 효과 +190~380%)", None),
    (1, "③ 썸네일·상세페이지·옵션 구성 개선", None),
], note="출처: comparePrice.html, 토스페이먼츠(2.06%), 한국소비자연맹(72.4%), Spiegel Research(190/380%)")

# ===== 10. 대책 P3 소싱 =====
content_slide("대책 P3", "카탈로그 확장·소싱 — 주문수 회복의 근본", [
    (0, "상품 수(노출 모수) 확대가 매출 회복의 근본 — 신상품 지속 등록", NAVY),
    (0, "2단계 소싱 전략", NAVY),
    (1, "1단계: 도매매(위탁판매) — MOQ 없음·무재고·발주 자동화로 상품수 빠르게 확대", None),
    (1, "2단계: 잘 팔리는 상품은 도매꾹 사입/브랜딩으로 전환 → 가격경쟁력·마진 확보", None),
    (1, "위탁은 경쟁 심해 마진 낮음 → 초기 확장용, 검증된 상품만 사입 전환", None),
    (0, "베스트/트렌드/시즌 상품 발굴 + 신상품 등록 주기 정례화", GREEN),
], note="출처: sweepingoms 소싱가이드, 토스페이먼츠 semo-29, tojobcn (3-0)")

# ===== 11. 어뷰징 금지 =====
content_slide("제약", "어뷰징 금지 — 매출전략의 절대 제약조건", [
    (0, "중복등록: 동일·근접 카테고리 동일상품 → 정렬점수 감점·판매제한·ID중지", RED),
    (0, "부정키워드(타 판매자명·무관 키워드): 1차 수정요청 → 2차 삭제·신용차감 → 3차 이용정지", RED),
    (0, "허위체결(자전거래): 계약해지·정산금 지급보류", RED),
    (0, "정책위반 상품은 노출순서 조정으로 상위 미노출 → 어뷰징은 매출 자충수", NAVY),
    (0, "합법적 랭크 개선 경로: 판매실적·리뷰·발송속도·CS(서비스품질) 축적", GREEN),
], note="출처: policy.gmarket.co.kr/terms-policy/seller-use, syi_policy.html (3-0)")

# ===== 12. 액션플랜 =====
table_slide("액션플랜", "우선순위 액션플랜 · KPI · 일정",
    ["우선","기간","핵심 액션","KPI"],
    [["P0","1주","14계정 판매·노출상태 전수점검→판매중 복구, 중복등록 정리","노출가능 상품수 회복(03수준)"],
     ["P1","2주","적자133 OFF→고ROAS51 증액, 무료배송·최저가, 검색용 상품명","적자광고비 0원화·노출/유입↑"],
     ["P2","1달","클릭0전환 가격·리뷰·상세 개선, 도매매 신상품 등록","전환율↑·신규 노출상품수↑"],
     ["상시","-","서비스품질(리뷰·발송·CS) 축적, 어뷰징 0건","랭크점수↑·패널티 0"]],
    widths=[1.2,1.3,6.5,4.3])

# ===== 13. 기대효과 =====
content_slide("기대효과", "기대 효과 (보수적 추정)", [
    (0, "즉효(P0+P1, ~2주): 적자광고비 월 50만원 낭비 차단 → 고ROAS로 재투입", GREEN),
    (1, "노출 정상화로 유입 회복 시 주문수 반등 여력 큼(03 대비 -65% 회복 목표)", None),
    (0, "단기(P2, ~1달): 전환율 개선 + 신상품 노출로 매출 우상향 전환", GREEN),
    (0, "정량 목표(예시): 노출가능 상품수 03수준 회복, ROAS 263%→300%+, 적자광고비 0원", NAVY),
    (1, "주의: 랭크 가중치·노출하락 원인은 비공개라 14계정 데이터로 지속 진단 필요", GRAY),
], note="수치 목표는 실데이터 기준 재보정 전제. 광고는 흑자라 '노출 정상화'가 매출 레버의 핵심")

# ===== 14. 출처 =====
content_slide("출처", "근거 출처 (공식 문서 중심)", [
    (0, "ESM PLUS 공식 판매자가이드 2.0", NAVY),
    (1, "comparePrice.html(노출순위·최저가) / goodsName.html(상품명) / goodsExposeWait.html(노출상태)", None),
    (1, "productRegisterIntegrate.html(상품2.0) / syi_policy.html(어뷰징정책)", None),
    (0, "ESM 마케팅허브·광고가이드", NAVY),
    (1, "marketinghub power-click/introduction(파워클릭) / ad-term / AD_Guide_all.pdf", None),
    (1, "policy.gmarket.co.kr/terms-policy/seller-use(판매자정책)", None),
    (0, "보조: i-boss(키워드), 토스페이먼츠(CVR/소싱), 한국소비자연맹·Spiegel(리뷰), sweepingoms(소싱)", GRAY),
    (1, "검증: 딥리서치 2회 · 32개 클레임 적대적 3표 검증(16+16 confirmed)", GRAY),
])

out = '/home/rejoice888/Avengers/frontend/public/지마켓_매출증가_대책보고서.pptx'
prs.save(out)
print('PPT 저장 완료:', out, '/ 슬라이드', len(prs.slides._sldIdLst))
